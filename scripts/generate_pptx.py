#!/usr/bin/env python3
"""
Standalone PPT generator for testing outside the Next.js pipeline.

Usage:
  python scripts/generate_pptx.py /path/to/workdir

Expects workdir to contain:
  - analysis_summary.json
  - chart_data.json
  - insights.md
  - (optional) client_logo.png
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    sys.exit(1)

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
except ImportError:
    print("Install matplotlib: pip install matplotlib")
    sys.exit(1)


# Default styling
COLORS = {
    "primary": RGBColor(0x1A, 0x73, 0xE8),
    "dark": RGBColor(0x20, 0x21, 0x24),
    "text": RGBColor(0x3C, 0x40, 0x43),
    "muted": RGBColor(0x80, 0x86, 0x8B),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

FONT = "Calibri"


def create_title_slide(prs: Presentation, title: str, subtitle: str, logo_path: str | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.5), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["dark"]
    p.font.name = FONT
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLORS["muted"]
    p2.font.name = FONT

    # Logo
    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.3), Inches(1.5))

    return slide


def create_insights_slide(prs: Presentation, insights: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Heading
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["dark"]
    p.font.name = FONT

    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.5), Inches(5.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, insight in enumerate(insights):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"• {insight}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["text"]
        p.font.name = FONT
        p.space_after = Pt(12)

    return slide


def create_metrics_slide(prs: Presentation, summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Heading
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Metrics"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["dark"]
    p.font.name = FONT

    # Display metrics as large numbers
    metrics = {k: v for k, v in summary.items() if k != "columns"}
    x_start = 0.75
    y_pos = 1.8
    col_width = 3.0

    for i, (key, value) in enumerate(metrics.items()):
        col = i % 4
        row = i // 4
        x = x_start + col * col_width
        y = y_pos + row * 2.0

        # Value
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.8), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.font.name = FONT

        # Label
        p2 = tf.add_paragraph()
        p2.text = key.replace("_", " ").title()
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLORS["muted"]
        p2.font.name = FONT

    return slide


def create_chart_slide(prs: Presentation, chart_data: list[dict], workdir: Path):
    if not chart_data:
        return None

    # Generate a bar chart with matplotlib
    labels = [d.get("label", d.get("gambit_id", "")) for d in chart_data[:15]]
    values = [d.get("count", 0) for d in chart_data[:15]]

    fig, ax = plt.subplots(figsize=(12, 5))
    bars = ax.barh(labels, values, color="#1a73e8", height=0.6)
    ax.set_xlabel("Count", fontsize=12)
    ax.set_title("Gambit Performance", fontsize=16, fontweight="bold", pad=15)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.invert_yaxis()
    plt.tight_layout()

    chart_path = workdir / "chart_gambit.png"
    fig.savefig(chart_path, dpi=150, bbox_inches="tight")
    plt.close()

    # Add to slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Gambit Analysis"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["dark"]
    p.font.name = FONT

    slide.shapes.add_picture(str(chart_path), Inches(0.75), Inches(1.5), Inches(11.5))

    return slide


def generate_report(workdir: str):
    workdir = Path(workdir)

    # Load inputs
    with open(workdir / "analysis_summary.json") as f:
        summary = json.load(f)

    with open(workdir / "chart_data.json") as f:
        chart_data = json.load(f)

    insights_path = workdir / "insights.md"
    insights = []
    if insights_path.exists():
        insights = [
            line.lstrip("- *").strip()
            for line in insights_path.read_text().splitlines()
            if line.strip().startswith(("-", "*"))
        ]

    logo_path = workdir / "client_logo.png"
    logo = str(logo_path) if logo_path.exists() else None

    # Build deck
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Bot Performance Report", "Automated Analysis", logo)
    create_insights_slide(prs, insights or ["Analysis completed."])
    create_metrics_slide(prs, summary)
    create_chart_slide(prs, chart_data, workdir)

    # Save
    output = workdir / "report.pptx"
    prs.save(str(output))
    print(f"Report saved to {output}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_pptx.py <workdir>")
        sys.exit(1)

    generate_report(sys.argv[1])
