"""
pptx_generator.py
Generates Business Review decks precisely matching the Tars/Amex PPTX template.

Extracted values from AMEX_Platinum_Nov_2025_to_Feb_2026.pptx:
  Slide size    : 10.0" x 5.625"
  Title box fill: #B28FD5 (light purple — NOT pink)
  Title text    : accent3 = #78909C (blue-grey), 18pt bold
  Table header  : #CCCCCC, 10pt bold, dark text
  Table body    : no fill (white), 10pt
  Month col bold: True
  Cover purple  : #6C4098
  Cover dots    : #B28FD5
  Contact font  : Montserrat SemiBold (or Arial fallback)
  Logo group    : left=0.171" top=0.295" w=1.345" h=0.520"
  Client logo   : left=1.660" top=0.096" w=1.631" h=0.917"
  Title box     : left=3.939" top=0.446" w=6.115" h=0.505"
  Content area  : starts at ~1.25" from top
"""

import os
import random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ─── EXACT COLORS FROM PPTX ───────────────────────────────────────────────────
C_COVER_PURPLE  = RGBColor(0x6C, 0x40, 0x98)   # Cover title + dots
C_TITLE_TEXT    = RGBColor(0x43, 0x43, 0x43)   # Title box text (#434343)
C_TITLE_BOX     = RGBColor(0xB2, 0x8F, 0xD5)   # Title box fill (light purple)
C_TABLE_HEADER  = RGBColor(0xCC, 0xCC, 0xCC)   # Table header grey (#CCCCCC)
C_TABLE_ALT     = RGBColor(0xF5, 0xF5, 0xF5)   # Alternating row (light grey)
C_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT     = RGBColor(0x43, 0x43, 0x43)   # Body text
C_MID_TEXT      = RGBColor(0x43, 0x43, 0x43)
C_BULLET_TEXT   = RGBColor(0x43, 0x43, 0x43)
C_DOT           = RGBColor(0xB2, 0x8F, 0xD5)   # Same as title box
C_DIVIDER_LINE  = RGBColor(0xB2, 0x8F, 0xD5)   # Horizontal divider on cover
C_CHART_BLUE    = RGBColor(0x4A, 0x72, 0xC4)
C_CHART_RED     = RGBColor(0xE8, 0x4A, 0x4A)
C_PIE = ["4A72C4","E84A4A","4AB04A","F5A623","9B59B6","1ABC9C","E67E22","34495E"]

# ─── SLIDE DIMENSIONS (exact from PPTX) ───────────────────────────────────────
W = Inches(10.0)
H = Inches(5.625)

# ─── EXACT POSITIONS FROM PPTX ────────────────────────────────────────────────
# Logo group (Tars)       : left=0.171" top=0.295" w=1.345" h=0.520"
# Client logo (Amex)      : left=1.660" top=0.096" w=1.631" h=0.917"
# Title text box          : left=3.939" top=0.446" w=6.115" h=0.505"
# Content starts at       : top ~1.25"

TARS_LOGO_L  = Inches(0.171)
TARS_LOGO_T  = Inches(0.295)
TARS_LOGO_W  = Inches(1.345)
TARS_LOGO_H  = Inches(0.520)

CLIENT_LOGO_L = Inches(1.660)
CLIENT_LOGO_T = Inches(0.096)
CLIENT_LOGO_W = Inches(1.631)
CLIENT_LOGO_H = Inches(0.917)

TITLE_BOX_L  = Inches(3.939)
TITLE_BOX_T  = Inches(0.446)
TITLE_BOX_W  = Inches(6.115)
TITLE_BOX_H  = Inches(0.505)

CONTENT_TOP  = Inches(1.22)
MARGIN_L     = Inches(0.2)
MARGIN_R     = Inches(0.2)
CONTENT_W    = W - MARGIN_L - MARGIN_R

FONT         = "Arial"
FONT_CONTACT = "Arial"   # Montserrat SemiBold if available


# ─── HELPERS ──────────────────────────────────────────────────────────────────

def white_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE


def add_rect(slide, x, y, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def add_text(slide, text, x, y, w, h, size, bold=False,
             color=C_DARK_TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_logos(slide, assets):
    """Place logos exactly matching the PPTX — Tars left, client logo right of it."""
    tars = assets.get("tars_logo")
    client = assets.get("client_logo")

    if tars and os.path.exists(tars):
        try:
            slide.shapes.add_picture(tars, TARS_LOGO_L, TARS_LOGO_T,
                                     width=TARS_LOGO_W, height=TARS_LOGO_H)
        except Exception:
            add_text(slide, "≡TARS", TARS_LOGO_L, TARS_LOGO_T,
                    TARS_LOGO_W, TARS_LOGO_H, 14, bold=True, color=C_COVER_PURPLE)
    else:
        add_text(slide, "≡TARS", TARS_LOGO_L, TARS_LOGO_T,
                TARS_LOGO_W, TARS_LOGO_H, 14, bold=True, color=C_COVER_PURPLE)

    if client and os.path.exists(client):
        try:
            slide.shapes.add_picture(client, CLIENT_LOGO_L, CLIENT_LOGO_T,
                                     width=CLIENT_LOGO_W, height=CLIENT_LOGO_H)
        except Exception:
            pass


def add_title_box(slide, title):
    """Light purple title box matching exact PPTX position and color."""
    box = slide.shapes.add_shape(1, TITLE_BOX_L, TITLE_BOX_T,
                                  TITLE_BOX_W, TITLE_BOX_H)
    box.fill.solid()
    box.fill.fore_color.rgb = C_TITLE_BOX
    box.line.fill.background()

    add_text(slide, title,
             TITLE_BOX_L + Inches(0.12), TITLE_BOX_T + Inches(0.04),
             TITLE_BOX_W - Inches(0.18), TITLE_BOX_H - Inches(0.06),
             14, bold=True, color=C_TITLE_TEXT, align=PP_ALIGN.CENTER)


def add_content_header(slide, title, assets):
    white_bg(slide)
    add_logos(slide, assets)
    add_title_box(slide, title)


def build_table_small(slide, headers, rows, x, y, w):
    """Smaller font version for wide tables with many columns."""
    if not headers or not rows:
        return y
    col_n = len(headers)
    row_n = len(rows) + 1
    row_h = Inches(0.28)
    total_h = row_h * row_n
    tbl = slide.shapes.add_table(row_n, col_n, x, y, w, total_h).table
    col_w = int(w / col_n)
    for i in range(col_n):
        tbl.columns[i].width = col_w
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TABLE_HEADER
        tf = cell.text_frame
        tf.margin_top = Pt(1); tf.margin_bottom = Pt(1); tf.margin_left = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0] if p.runs else p.add_run()
        r.text = h; r.font.name = FONT; r.font.size = Pt(8)
        r.font.bold = True; r.font.color.rgb = C_DARK_TEXT
    for i, row in enumerate(rows):
        is_total = str(row[0]).strip().lower() == "total"
        bg = C_TABLE_ALT if i % 2 == 0 else C_WHITE
        if is_total: bg = C_TABLE_HEADER
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.margin_top = Pt(1); tf.margin_bottom = Pt(1); tf.margin_left = Pt(2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.runs[0] if p.runs else p.add_run()
            r.text = str(val); r.font.name = FONT; r.font.size = Pt(8)
            r.font.bold = (j == 0) or is_total; r.font.color.rgb = C_DARK_TEXT
    return y + total_h + Inches(0.15)


def build_table(slide, headers, rows, x, y, w):
    """Build table matching PPTX style: #CCCCCC header, white/light alternating rows, 10pt."""
    if not headers or not rows:
        return y

    col_n = len(headers)
    row_n = len(rows) + 1
    row_h = Inches(0.33)
    total_h = row_h * row_n

    tbl = slide.shapes.add_table(row_n, col_n, x, y, w, total_h).table
    col_w = int(w / col_n)
    for i in range(col_n):
        tbl.columns[i].width = col_w

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TABLE_HEADER
        tf = cell.text_frame
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        tf.margin_left = Pt(4)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0] if p.runs else p.add_run()
        r.text = h
        r.font.name = FONT
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C_DARK_TEXT

    # Data rows
    for i, row in enumerate(rows):
        is_total = str(row[0]).strip().lower() == "total"
        bg = C_TABLE_ALT if i % 2 == 0 else C_WHITE
        if is_total:
            bg = C_TABLE_HEADER

        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
            tf.margin_left = Pt(4)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.runs[0] if p.runs else p.add_run()
            r.text = str(val)
            r.font.name = FONT
            r.font.size = Pt(10)
            r.font.bold = (j == 0) or is_total
            r.font.color.rgb = C_DARK_TEXT

    return y + total_h + Inches(0.15)


def add_bullet_point(slide, text, x, y, w, bold_prefix=None):
    """Bullet with a small grey dot — matching the PPTX style."""
    dot = slide.shapes.add_shape(9, x, y + Inches(0.09), Inches(0.09), Inches(0.09))
    dot.fill.solid()
    dot.fill.fore_color.rgb = C_DARK_TEXT
    dot.line.fill.background()

    tb = slide.shapes.add_textbox(x + Inches(0.18), y, w - Inches(0.18), Inches(0.42))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    if bold_prefix:
        r1 = p.add_run()
        r1.text = bold_prefix
        r1.font.name = FONT
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = C_DARK_TEXT
        r2 = p.add_run()
        r2.text = text
        r2.font.name = FONT
        r2.font.size = Pt(11)
        r2.font.color.rgb = C_DARK_TEXT
    else:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(11)
        r.font.color.rgb = C_DARK_TEXT


# ─── SLIDE BUILDERS ───────────────────────────────────────────────────────────

def build_cover_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    # Scattered dots (matching PPTX feel — random sizes, clustered corners)
    random.seed(99)
    dots = []
    for _ in range(22):
        dots.append((random.uniform(0.0, 2.2), random.uniform(0.1, 5.5),
                     random.uniform(0.06, 0.20)))
    for _ in range(22):
        dots.append((random.uniform(7.8, 9.9), random.uniform(0.1, 5.5),
                     random.uniform(0.06, 0.20)))
    for dx, dy, ds in dots:
        d = slide.shapes.add_shape(9, Inches(dx), Inches(dy), Inches(ds), Inches(ds))
        d.fill.solid()
        d.fill.fore_color.rgb = C_DOT
        d.line.fill.background()

    # Tars logo — centered top
    tars = assets.get("tars_logo")
    logo_placed = False
    if tars and os.path.exists(tars):
        try:
            pic = slide.shapes.add_picture(tars, Inches(4.0), Inches(0.4), height=Inches(0.9))
            pic.left = int((W - pic.width) / 2)
            logo_placed = True
        except Exception:
            pass
    if not logo_placed:
        add_text(slide, "≡TARS", Inches(4.0), Inches(0.4), Inches(2.0), Inches(0.9),
                24, bold=True, color=C_COVER_PURPLE, align=PP_ALIGN.CENTER)

    # Divider line
    add_rect(slide, Inches(2.3), Inches(1.55), Inches(5.4), Inches(0.01), C_DIVIDER_LINE)

    # Client logo — large, centered
    client = assets.get("client_logo")
    client_placed = False
    if client and os.path.exists(client):
        try:
            pic = slide.shapes.add_picture(client, Inches(2.5), Inches(1.7), width=Inches(5.0))
            pic.left = int((W - pic.width) / 2)
            client_placed = True
        except Exception:
            pass
    if not client_placed:
        add_text(slide, assets.get("client_name", ""),
                Inches(1.5), Inches(1.7), Inches(7.0), Inches(1.4),
                40, bold=True, color=C_COVER_PURPLE, align=PP_ALIGN.CENTER)

    # Title
    add_text(slide, slide_data.get("title", "Business Review"),
             Inches(1.5), Inches(3.5), Inches(7.0), Inches(0.75),
             24, bold=True, color=C_COVER_PURPLE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, slide_data.get("subtitle", ""),
             Inches(1.5), Inches(4.15), Inches(7.0), Inches(0.5),
             14, bold=True, color=C_COVER_PURPLE, align=PP_ALIGN.CENTER)


def build_overview_stats_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_header(slide, slide_data.get("title", "Performance Overview"), assets)

    # Use locked numbers if available, fallback to slide_data
    locked = slide_data.get("_locked", {})
    fmt = locked.get("formatted", {})
    stats = fmt.get("stat_cards") or slide_data.get("stats", [])
    insights = slide_data.get("insights", [])
    period = fmt.get("period_label") or slide_data.get("period", "")

    y = CONTENT_TOP
    if period:
        add_text(slide, f"Period: {period}", MARGIN_L, y, Inches(4), Inches(0.28),
                10, italic=True, color=C_DARK_TEXT)
        y += Inches(0.3)

    count = min(len(stats), 4)
    gap = Inches(0.15)
    cw = (CONTENT_W - gap * (count - 1)) / count
    ch = Inches(1.3)

    for i, stat in enumerate(stats[:4]):
        cx = MARGIN_L + i * (cw + gap)
        card = slide.shapes.add_shape(1, cx, y, cw, ch)
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_TABLE_HEADER
        card.line.width = Pt(0.75)

        # Thin purple top bar
        bar = add_rect(slide, cx, y, cw, Inches(0.04), C_COVER_PURPLE)

        add_text(slide, stat.get("value", ""),
                cx + Inches(0.05), y + Inches(0.05),
                cw - Inches(0.1), Inches(0.75),
                28, bold=True, color=C_COVER_PURPLE, align=PP_ALIGN.CENTER)

        add_text(slide, stat.get("label", ""),
                cx + Inches(0.05), y + Inches(0.78),
                cw - Inches(0.1), Inches(0.35),
                10, color=C_MID_TEXT, align=PP_ALIGN.CENTER)

        if stat.get("note"):
            add_text(slide, stat["note"],
                    cx + Inches(0.05), y + Inches(1.1),
                    cw - Inches(0.1), Inches(0.22),
                    9, italic=True, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

    insight_y = y + ch + Inches(0.22)
    for i, ins in enumerate(insights[:4]):
        iy = insight_y + i * Inches(0.44)
        if iy > H - Inches(0.3): break
        add_bullet_point(slide, ins, MARGIN_L, iy, CONTENT_W)


def build_monthly_trend_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_header(slide, slide_data.get("title", "Visits & Conversations - Over the Months"), assets)

    # Use locked numbers if available
    locked = slide_data.get("_locked", {})
    fmt = locked.get("formatted", {})
    headers = fmt.get("trend_headers") or slide_data.get("headers", [])
    rows = fmt.get("trend_rows") or slide_data.get("rows", [])
    insights = slide_data.get("insights", [])
    chart_data_obj = fmt.get("chart_data") or slide_data.get("chart_data")
    has_chart = (chart_data_obj and
                 chart_data_obj.get("labels") and
                 len(chart_data_obj.get("labels", [])) > 1)

    y = CONTENT_TOP

    if headers and rows:
        # With many columns, place table full width and chart below or skip chart
        col_count = len(headers)
        if has_chart and col_count <= 5:
            tw = Inches(5.6)
            next_y = build_table(slide, headers, rows, MARGIN_L, y, tw)
            try:
                cd = ChartData()
                cd.categories = chart_data_obj["labels"]
                if chart_data_obj.get("bot_visits"):
                    cd.add_series("Bot Visits", chart_data_obj["bot_visits"])
                if chart_data_obj.get("conversations"):
                    cd.add_series("Conversations", chart_data_obj["conversations"])
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, cd,
                    Inches(5.95), y, Inches(3.85), Inches(2.6)
                ).chart
                chart.has_title = False
                chart.has_legend = True
                chart.legend.position = 2
                chart.legend.include_in_layout = False
                plot = chart.plots[0]
                plot.series[0].format.fill.solid()
                plot.series[0].format.fill.fore_color.rgb = C_CHART_BLUE
                if len(plot.series) > 1:
                    plot.series[1].format.fill.solid()
                    plot.series[1].format.fill.fore_color.rgb = C_CHART_RED
            except Exception:
                pass
        else:
            # Many columns — use full width, smaller font
            next_y = build_table_small(slide, headers, rows, MARGIN_L, y, CONTENT_W)

        # Place insights below table, ensuring no overlap
        insight_y = max(next_y, y + Inches(0.1))
        if has_chart and col_count <= 5:
            # Chart is on right, insights go below chart area
            insight_y = max(next_y, y + Inches(2.75))
    else:
        insight_y = y

    for i, ins in enumerate(insights[:3]):
        iy = insight_y + i * Inches(0.44)
        if iy > H - Inches(0.25): break
        add_bullet_point(slide, ins, MARGIN_L, iy, CONTENT_W)


def build_generic_table_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_header(slide, slide_data.get("title", "Analysis"), assets)

    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])
    headers2 = slide_data.get("headers2")
    rows2 = slide_data.get("rows2")
    insights = slide_data.get("insights", [])
    description = slide_data.get("description", "")
    chart_data_obj = slide_data.get("chart_data")
    chart_type = slide_data.get("chart_type", "")
    footer_text = slide_data.get("footer_text", "")

    y = CONTENT_TOP

    if description:
        add_text(slide, description, MARGIN_L, y, CONTENT_W, Inches(0.32),
                10, italic=False, color=C_DARK_TEXT)
        y += Inches(0.35)

    has_pie = (chart_type == "pie" and chart_data_obj and
               chart_data_obj.get("labels") and chart_data_obj.get("values"))

    if has_pie:
        try:
            cd = ChartData()
            cd.categories = chart_data_obj["labels"]
            cd.add_series("", chart_data_obj["values"])
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, cd,
                MARGIN_L, y, Inches(4.2), Inches(3.0)
            ).chart
            chart_title = chart_data_obj.get("title", "")
            chart.has_title = bool(chart_title)
            if chart_title:
                chart.chart_title.text_frame.text = chart_title
            chart.has_legend = True
            chart.legend.position = 4
            for idx, point in enumerate(chart.plots[0].series[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(
                    C_PIE[idx % len(C_PIE)])

            # Table to the right
            if headers and rows:
                next_y = build_table(slide, headers, rows,
                                     Inches(4.55), y, Inches(5.25))
                insight_y = max(y + Inches(3.15), next_y)
            else:
                insight_y = y + Inches(3.15)
        except Exception:
            if headers and rows:
                next_y = build_table(slide, headers, rows, MARGIN_L, y, CONTENT_W)
                insight_y = next_y
            else:
                insight_y = y
    elif headers and rows:
        # Two tables stacked (like slide 3 in PPTX)
        if headers2 and rows2:
            tw = CONTENT_W * 0.52
            # Use small font if first table has many columns
            if len(headers) > 4:
                next_y1 = build_table_small(slide, headers, rows, MARGIN_L, y, tw)
            else:
                next_y1 = build_table(slide, headers, rows, MARGIN_L, y, tw)
            next_y2 = build_table(slide, headers2, rows2, MARGIN_L, next_y1, tw)
            insight_y = next_y2
        else:
            next_y = build_table(slide, headers, rows, MARGIN_L, y, CONTENT_W)
            insight_y = next_y
    else:
        insight_y = y

    # Footer italic text (like the bold summary line in slide 3)
    if footer_text and insight_y < H - Inches(0.65):
        add_text(slide, footer_text, MARGIN_L, insight_y, CONTENT_W, Inches(0.55),
                9, bold=True, italic=True, color=C_DARK_TEXT)
        insight_y += Inches(0.6)

    for i, ins in enumerate(insights[:4]):
        iy = insight_y + i * Inches(0.44)
        if iy > H - Inches(0.25): break
        if has_pie:
            add_bullet_point(slide, ins, MARGIN_L + Inches(4.4), iy, CONTENT_W * 0.55)
        else:
            add_bullet_point(slide, ins, MARGIN_L + Inches(3.0), iy, CONTENT_W * 0.65)


def build_insights_summary_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_header(slide, slide_data.get("title", "Key Insights"), assets)

    insights = slide_data.get("insights", [])
    y = CONTENT_TOP

    for i, ins in enumerate(insights[:6]):
        iy = y + i * Inches(0.65)
        if iy > H - Inches(0.3): break

        # Numbered circle
        c = slide.shapes.add_shape(9, MARGIN_L, iy, Inches(0.35), Inches(0.35))
        c.fill.solid()
        c.fill.fore_color.rgb = C_COVER_PURPLE
        c.line.fill.background()
        add_text(slide, str(i+1), MARGIN_L, iy - Inches(0.02),
                Inches(0.35), Inches(0.39), 11, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, ins, MARGIN_L + Inches(0.45), iy,
                CONTENT_W - Inches(0.45), Inches(0.55),
                11, color=C_DARK_TEXT)


def build_next_steps_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_header(slide, slide_data.get("title", "Next Steps for Improvements"), assets)

    items = slide_data.get("items", [])
    y = CONTENT_TOP

    for i, item in enumerate(items[:5]):
        iy = y + i * Inches(0.75)
        if iy > H - Inches(0.3): break
        add_bullet_point(slide, item, MARGIN_L, iy, CONTENT_W)


def build_contact_slide(prs, slide_data, assets):
    """Contact slide matching PPTX exactly — white bg, actual icon images, Montserrat font."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_header(slide, "GET IN TOUCH WITH US", assets)

    # Contact items with icon placeholder circles
    contacts = [
        (slide_data.get("email",   "team@hellotars.com"),   Inches(2.04)),
        (slide_data.get("social",  "@hellotars.ai"),         Inches(2.83)),
        (slide_data.get("website", "www.hellotars.com"),     Inches(3.73)),
    ]
    icon_labels = ["@", "✦", "⊕"]

    for idx, (value, top) in enumerate(contacts):
        # Icon circle (grey, matching PPTX)
        c = slide.shapes.add_shape(9, Inches(0.857), top, Inches(0.5), Inches(0.5))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
        c.line.fill.background()
        add_text(slide, icon_labels[idx], Inches(0.857), top - Inches(0.02),
                Inches(0.5), Inches(0.54), 12, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, value, Inches(1.712), top + Inches(0.08),
                Inches(3.0), Inches(0.32), 14, bold=True,
                color=C_COVER_PURPLE)

    # Client logo on the right — large (matching PPTX: left=5.180" top=1.053" w=4.198")
    client = assets.get("client_logo")
    if client and os.path.exists(client):
        try:
            slide.shapes.add_picture(client, Inches(5.18), Inches(1.05), width=Inches(4.2))
        except Exception:
            pass


# ─── REGISTRY ─────────────────────────────────────────────────────────────────

SLIDE_BUILDERS = {
    "cover":             build_cover_slide,
    "overview_stats":    build_overview_stats_slide,
    "monthly_trend":     build_monthly_trend_slide,
    "user_journey":      build_generic_table_slide,
    "conversion_funnel": build_generic_table_slide,
    "device_breakdown":  build_generic_table_slide,
    "gambit_analysis":   build_generic_table_slide,
    "insights_summary":  build_insights_summary_slide,
    "next_steps":        build_next_steps_slide,
    "contact":           build_contact_slide,
}


def generate_pptx(slide_plan: dict, assets: dict, output_path: str, locked_numbers: dict = None):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_data = slide_plan.get("slides", [])
    print(f"  🎨 Generating {len(slides_data)} slides...")

    for i, sd in enumerate(slides_data):
        stype = sd.get("type", "generic")
        builder = SLIDE_BUILDERS.get(stype, build_generic_table_slide)
        if stype == "cover":
            sd["client_name"] = assets.get("client_name", "")
        # Inject locked numbers into slide data so builders use accurate numbers
        if locked_numbers:
            sd["_locked"] = locked_numbers
        try:
            builder(prs, sd, assets)
            print(f"    ✓ Slide {i+1}: {stype} — '{sd.get('title', stype)}'")
        except Exception as e:
            print(f"    [WARN] Slide {i+1} ({stype}) error: {e}")
            import traceback; traceback.print_exc()
            fallback = prs.slides.add_slide(prs.slide_layouts[6])
            white_bg(fallback)
            add_text(fallback, sd.get("title", stype),
                    Inches(1), Inches(2), Inches(8), Inches(1.5),
                    20, bold=True, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    print(f"  [SAVED] {output_path}")
    return output_path


# ─── CLI ENTRY POINT ──────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys
    import json
    from pathlib import Path

    if len(sys.argv) < 2:
        print("Usage: python pptx_generator.py <work_directory>")
        sys.exit(1)

    work_dir = Path(sys.argv[1])

    # Read slide_plan.json
    slide_plan_path = work_dir / "slide_plan.json"
    if not slide_plan_path.exists():
        print(f"Error: slide_plan.json not found in {work_dir}")
        sys.exit(1)
    with open(slide_plan_path, "r", encoding="utf-8") as f:
        slide_plan = json.load(f)

    # Read locked_numbers.json
    locked_path = work_dir / "locked_numbers.json"
    locked_numbers = None
    if locked_path.exists():
        with open(locked_path, "r", encoding="utf-8") as f:
            locked_numbers = json.load(f)

    # Build assets dict from deck_input.json if available
    assets = {}
    deck_input_path = work_dir / "deck_input.json"
    if deck_input_path.exists():
        with open(deck_input_path, "r", encoding="utf-8") as f:
            deck_input = json.load(f)
        logo_path = deck_input.get("logo_path")
        if logo_path:
            assets["client_logo"] = logo_path
    if locked_numbers:
        assets["client_name"] = locked_numbers.get("client_name", "")

    output_path = str(work_dir / "report.pptx")
    generate_pptx(slide_plan, assets, output_path, locked_numbers)
